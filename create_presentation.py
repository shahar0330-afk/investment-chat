#!/usr/bin/env python3
"""Generate investment chat business plan presentation — v2 with RTL fix, images, and better design."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

ASSETS = os.path.expanduser('~/investment-chat/assets')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
BG_DARK = RGBColor(0x0b, 0x0b, 0x1e)
BG_CARD = RGBColor(0x15, 0x15, 0x2b)
BG_CARD_ALT = RGBColor(0x1a, 0x1a, 0x35)
ACCENT = RGBColor(0x6c, 0x63, 0xff)
ACCENT_LIGHT = RGBColor(0x9b, 0x93, 0xff)
ACCENT_DIM = RGBColor(0x3d, 0x37, 0x99)
GOLD = RGBColor(0xfb, 0xbf, 0x24)
GOLD_DIM = RGBColor(0x92, 0x6f, 0x13)
GREEN = RGBColor(0x4a, 0xde, 0x80)
GREEN_DIM = RGBColor(0x16, 0x65, 0x34)
RED = RGBColor(0xf8, 0x71, 0x71)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x88, 0x88, 0xa0)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xdd)
DARK_OVERLAY = RGBColor(0x08, 0x08, 0x18)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_rtl(paragraph):
    """Set RTL direction on a paragraph via XML."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('rtl', '1')


def add_shape(slide, left, top, width, height, fill_color=BG_CARD, corner_radius=None, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if alpha is not None:
        # Set transparency via XML on the spPr element
        spPr = shape._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', str(int(alpha * 1000)))
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.RIGHT, font_name='Arial', rtl=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if rtl:
        set_rtl(p)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT, spacing=8, rtl=True):
    """Add text with multiple paragraphs, each properly RTL."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.alignment = alignment
        p.space_after = Pt(spacing)
        if rtl:
            set_rtl(p)
    return txBox


def add_image_with_overlay(slide, img_path, left, top, width, height, overlay_alpha=70):
    """Add image with a dark overlay shape on top."""
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)
        overlay = add_shape(slide, left, top, width, height, DARK_OVERLAY, alpha=overlay_alpha)
        return True
    return False


def add_accent_line(slide, left, top, width, color=ACCENT):
    add_shape(slide, left, top, width, Inches(0.05), color)


def add_glow_circle(slide, cx, cy, radius, color):
    """Add a glowing circle accent."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        cx - radius, cy - radius, radius * 2, radius * 2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════
# SLIDE 1: Title — with background image
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_image_with_overlay(slide, f'{ASSETS}/finance-bg.jpg',
                       Inches(0), Inches(0), Inches(13.333), Inches(7.5), overlay_alpha=75)

# Decorative circles
add_glow_circle(slide, Inches(1.5), Inches(1.5), Inches(0.4), ACCENT_DIM)
add_glow_circle(slide, Inches(11.5), Inches(6), Inches(0.3), GOLD_DIM)

# Accent bars
add_shape(slide, Inches(4.5), Inches(3.05), Inches(4.333), Inches(0.06), ACCENT)

add_text(slide, Inches(1), Inches(1.2), Inches(11.333), Inches(1.4),
         'נועם', font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.3), Inches(11.333), Inches(0.8),
         'יועץ השקעות AI', font_size=40, color=ACCENT_LIGHT, bold=False,
         alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(3.5), Inches(9.333), Inches(0.8),
         'צ׳אט חכם שמבין אותך, מנתח את השוק, ובונה לך תיק השקעות מותאם אישית',
         font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

# iPhone mockup on right side
iphone_chat = f'{ASSETS}/iphone-screen-chat.png'
if os.path.exists(iphone_chat):
    slide.shapes.add_picture(iphone_chat, Inches(9.8), Inches(1.5), Inches(2.5))

# Bottom card
card = add_shape(slide, Inches(2), Inches(5), Inches(6.333), Inches(1.5), BG_CARD, 0.08, alpha=85)
add_text(slide, Inches(2), Inches(5.15), Inches(6.333), Inches(0.5),
         'תכנון עסקי  |  2026', font_size=20, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(5.7), Inches(6.333), Inches(0.5),
         'investment-chat-psi.vercel.app', font_size=16, color=ACCENT_LIGHT,
         alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 2: הבעיה — with chart image
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Right side image
add_image_with_overlay(slide, f'{ASSETS}/chart.jpg',
                       Inches(8.5), Inches(0), Inches(4.833), Inches(7.5), overlay_alpha=60)

add_text(slide, Inches(0.8), Inches(0.4), Inches(7), Inches(0.9),
         'הבעיה', font_size=44, color=GOLD, bold=True)

add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(1.5), GOLD)

add_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(0.6),
         'ייעוץ השקעות בישראל שבור — ואף אחד לא פותר את זה', font_size=18, color=GRAY)

problems = [
    ('💸', 'ייעוץ השקעות יקר מדי',
     'יועץ אנושי גובה 0.5%-1% מהתיק בשנה.\nעל תיק של 500K₪ זה עד 5,000₪ בשנה.'),
    ('🚫', 'חוסר נגישות',
     '80% מהישראלים לא מקבלים ייעוץ מקצועי.\nמשקיעים לפי תחושה או לא משקיעים בכלל.'),
    ('📦', 'מוצרים גנריים',
     'הבנקים מציעים את אותם מוצרים לכולם.\nאין התאמה אמיתית לפרופיל ואמונות הלקוח.'),
    ('⏳', 'תהליך ארוך ומייגע',
     'פגישות, טפסים, זמני המתנה ארוכים.\nהדור הצעיר רוצה פתרון דיגיטלי מיידי.'),
]

for i, (icon, title, desc) in enumerate(problems):
    y = Inches(2.2 + i * 1.3)
    # Left accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.15), Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()

    add_text(slide, Inches(1.2), y - Inches(0.05), Inches(0.6), Inches(0.6),
             icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.8), y, Inches(5.5), Inches(0.4),
             title, font_size=20, color=WHITE, bold=True)
    add_multiline_text(slide, Inches(1.8), y + Inches(0.4), Inches(5.5), Inches(0.8),
                       desc.split('\n'), font_size=13, color=GRAY, spacing=2)


# ═══════════════════════════════════════════════════════
# SLIDE 3: הפתרון — with dashboard image
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# iPhone mockup on left — portfolio view
iphone_portfolio = f'{ASSETS}/iphone-screen-portfolio.png'
if os.path.exists(iphone_portfolio):
    slide.shapes.add_picture(iphone_portfolio, Inches(1), Inches(0.8), Inches(3.2))

add_text(slide, Inches(5.5), Inches(0.4), Inches(7.5), Inches(0.9),
         'הפתרון: נועם', font_size=44, color=ACCENT, bold=True)

add_accent_line(slide, Inches(5.5), Inches(1.2), Inches(1.5), ACCENT)

add_text(slide, Inches(5.5), Inches(1.5), Inches(7.3), Inches(0.6),
         'יועץ השקעות AI שמדבר עברית ומבין אותך באמת', font_size=18, color=GRAY)

features = [
    ('💬', 'שיחה טבעית בעברית',
     'לא שאלון — שיחה אמיתית. שואל שאלות חכמות ומבין את הלקוח לעומק.'),
    ('🧠', 'מבין את הפרופיל שלך',
     'מזהה רמת סיכון, אמונות בסקטורים, אופק השקעה, וניסיון קודם.'),
    ('📊', 'נתונים אמיתיים בזמן אמת',
     'מושך מידע חי: מכפילים, רווחיות, מצב שוק, מדדים, ואג"ח.'),
    ('📋', 'תיק מותאם אישית',
     'בונה הקצאת נכסים מלאה: מניות + אג"ח + מדדים + סכומים.'),
]

for i, (icon, title, desc) in enumerate(features):
    y = Inches(2.2 + i * 1.3)
    card = add_shape(slide, Inches(5.5), y, Inches(7.3), Inches(1.1), BG_CARD, 0.05, alpha=90)

    add_text(slide, Inches(5.8), y + Inches(0.1), Inches(0.6), Inches(0.5),
             icon, font_size=26, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(6.5), y + Inches(0.08), Inches(5.8), Inches(0.4),
             title, font_size=19, color=WHITE, bold=True)
    add_text(slide, Inches(6.5), y + Inches(0.5), Inches(5.8), Inches(0.5),
             desc, font_size=13, color=GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 3.5: Demo — dual iPhone screens
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_image_with_overlay(slide, f'{ASSETS}/finance-bg.jpg',
                       Inches(0), Inches(0), Inches(13.333), Inches(7.5), overlay_alpha=85)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9),
         'ככה זה נראה', font_size=44, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.8), Inches(1.1), Inches(1.8), GOLD)

# Chat iPhone on right
iphone_chat_demo = f'{ASSETS}/iphone-screen-chat.png'
if os.path.exists(iphone_chat_demo):
    slide.shapes.add_picture(iphone_chat_demo, Inches(1.5), Inches(1.4), Inches(3.3))

# Label under chat
add_text(slide, Inches(1.5), Inches(7), Inches(3.3), Inches(0.4),
         'שיחת הכרות עם היועץ', font_size=16, color=ACCENT_LIGHT, bold=True,
         alignment=PP_ALIGN.CENTER)

# Arrow between phones
add_text(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.8),
         '→', font_size=60, color=GOLD, alignment=PP_ALIGN.CENTER, rtl=False)
add_text(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.6),
         'בניית פרופיל\nוניתוח שוק', font_size=14, color=GRAY,
         alignment=PP_ALIGN.CENTER)

# Portfolio iPhone on left
iphone_portfolio_demo = f'{ASSETS}/iphone-screen-portfolio.png'
if os.path.exists(iphone_portfolio_demo):
    slide.shapes.add_picture(iphone_portfolio_demo, Inches(8.5), Inches(1.4), Inches(3.3))

# Label under portfolio
add_text(slide, Inches(8.5), Inches(7), Inches(3.3), Inches(0.4),
         'תיק השקעות מותאם אישית', font_size=16, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 4: איך זה עובד — flow with arrows
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'איך זה עובד?', font_size=44, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.8), Inches(1.2), Inches(1.8), GOLD)

# Connection line
add_shape(slide, Inches(1.8), Inches(3.95), Inches(9.8), Inches(0.04), ACCENT_DIM)

steps = [
    ('01', 'הכרות', 'הלקוח נכנס לצ׳אט ונועם\nמתחיל שיחה טבעית\nלהכיר אותו', '💬', ACCENT),
    ('02', 'פרופיל', 'מזהה רמת סיכון,\nסקטורים מועדפים,\nאופק וסכום השקעה', '🎯', GREEN),
    ('03', 'ניתוח שוק', 'AI מושך נתונים חיים\nעל מניות, מדדים,\nאג"ח ומצב השוק', '📊', GOLD),
    ('04', 'תיק מותאם', 'הלקוח מקבל תיק מלא:\nמניות, אג"ח, מדדים,\nסכומים ונימוקים', '📋', ACCENT_LIGHT),
]

for i, (num, title, desc, icon, color) in enumerate(steps):
    x = Inches(0.4 + i * 3.2)

    # Circle on line
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), Inches(3.72), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, x + Inches(1.1), Inches(3.72), Inches(0.5), Inches(0.5),
             icon, font_size=18, alignment=PP_ALIGN.CENTER, rtl=False)

    # Number above
    add_text(slide, x, Inches(1.7), Inches(2.8), Inches(0.6),
             num, font_size=48, color=color, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Arial', rtl=False)
    add_text(slide, x, Inches(2.4), Inches(2.8), Inches(0.5),
             title, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Card below
    card = add_shape(slide, x + Inches(0.1), Inches(4.5), Inches(2.6), Inches(2.3), BG_CARD, 0.06)
    add_multiline_text(slide, x + Inches(0.3), Inches(4.7), Inches(2.2), Inches(1.8),
                       desc.split('\n'), font_size=14, color=GRAY,
                       alignment=PP_ALIGN.CENTER, spacing=6)


# ═══════════════════════════════════════════════════════
# SLIDE 5: שוק היעד — with target image
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_image_with_overlay(slide, f'{ASSETS}/target.jpg',
                       Inches(0), Inches(0), Inches(13.333), Inches(7.5), overlay_alpha=82)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'שוק היעד', font_size=44, color=ACCENT, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.8), Inches(1.2), Inches(1.8), ACCENT)

# TAM SAM SOM — concentric feel
markets = [
    ('TAM', 'שוק עולמי', '$30B', 'Robo Advisory עד 2028', GREEN, Inches(3.5)),
    ('SAM', 'שוק ישראלי', '2B₪', 'ייעוץ השקעות דיגיטלי', GOLD, Inches(3)),
    ('SOM', 'יעד ראשוני', '50M₪', 'משקיעים צעירים 25-45', ACCENT_LIGHT, Inches(2.5)),
]

for i, (label, title, size, desc, color, card_w) in enumerate(markets):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(2.6), BG_CARD, 0.06, alpha=90)

    # Color top border
    add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(0.06), color)

    add_text(slide, x, Inches(2), Inches(3.8), Inches(0.45),
             label, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER,
             rtl=False)
    add_text(slide, x, Inches(2.5), Inches(3.8), Inches(0.7),
             size, font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             rtl=False)
    add_text(slide, x, Inches(3.15), Inches(3.8), Inches(0.35),
             title, font_size=14, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.5), Inches(3.8), Inches(0.5),
             desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Target segments
add_text(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.5),
         'קהלי יעד', font_size=26, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

segments = [
    ('🧑‍💻', 'צעירים 25-35', 'טק-סאבי, תיק השקעות ראשון, מחפשים כלי חכם'),
    ('👨‍👩‍👧', 'משפחות 35-50', 'חיסכון לילדים ולפנסיה, רוצים גישה מקצועית'),
    ('🏦', 'משקיעים מנוסים', 'רוצים AI כ-second opinion ולניתוח מהיר'),
    ('🌍', 'ישראלים בחו"ל', 'רוצים ייעוץ בעברית על שוק ישראלי ואמריקאי'),
]

for i, (icon, title, desc) in enumerate(segments):
    x = Inches(0.5 + i * 3.2)
    card = add_shape(slide, x, Inches(5.6), Inches(3), Inches(1.5), BG_CARD, 0.06, alpha=90)
    add_text(slide, x + Inches(0.15), Inches(5.7), Inches(0.65), Inches(0.6),
             icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.8), Inches(5.7), Inches(2), Inches(0.35),
             title, font_size=15, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.8), Inches(6.1), Inches(2), Inches(0.6),
             desc, font_size=11, color=GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 6: מודל עסקי — pricing cards
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'מודל עסקי', font_size=44, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.8), Inches(1.2), Inches(1.8), GOLD)

tiers = [
    ('חינם', '0₪', ['3 שיחות בחודש', 'תיק בסיסי', 'מניות בלבד', 'ללא נתונים חיים'], GRAY, BG_CARD),
    ('Pro', '49₪/חודש', ['שיחות ללא הגבלה', 'תיק מלא: מניות+אג"ח+מדדים', 'נתונים בזמן אמת', 'עדכון תיק חודשי'], ACCENT, BG_CARD_ALT),
    ('Premium', '149₪/חודש', ['הכל ב-Pro בתוספת:', 'ניתוח טכני מעמיק', 'התראות שוק אישיות', 'שיחה עם יועץ אנושי'], GOLD, BG_CARD),
]

for i, (name, price, features, color, bg) in enumerate(tiers):
    x = Inches(0.8 + i * 4.1)
    is_pro = name == 'Pro'
    y_offset = Inches(-0.15) if is_pro else Inches(0)
    h_extra = Inches(0.3) if is_pro else Inches(0)

    card = add_shape(slide, x, Inches(1.7) + y_offset, Inches(3.8), Inches(4.8) + h_extra, bg, 0.06)

    # Top color bar
    add_shape(slide, x, Inches(1.7) + y_offset, Inches(3.8), Inches(0.08), color)

    if is_pro:
        # "Most popular" badge
        badge = add_shape(slide, x + Inches(0.9), Inches(1.3), Inches(2), Inches(0.35), ACCENT, 0.1)
        add_text(slide, x + Inches(0.9), Inches(1.3), Inches(2), Inches(0.35),
                 'הכי פופולרי', font_size=11, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text(slide, x, Inches(2.05) + y_offset, Inches(3.8), Inches(0.45),
             name, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER,
             rtl=False)
    add_text(slide, x, Inches(2.55) + y_offset, Inches(3.8), Inches(0.6),
             price, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Separator
    add_shape(slide, x + Inches(0.5), Inches(3.2) + y_offset, Inches(2.8), Inches(0.02), GRAY)

    for j, feat in enumerate(features):
        check_color = GREEN if is_pro else (GOLD if name == 'Premium' else GRAY)
        add_text(slide, x + Inches(0.4), Inches(3.45 + j * 0.5) + y_offset, Inches(3.2), Inches(0.4),
                 f'✓  {feat}', font_size=14, color=LIGHT_GRAY)

# Bottom metrics
metrics_card = add_shape(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.6), BG_CARD, 0.05)
metrics = ['LTV צפוי: 1,200₪', 'CAC יעד: 150₪', 'LTV:CAC = 8:1', 'Payback: 3 חודשים']
for i, m in enumerate(metrics):
    x = Inches(1 + i * 3)
    add_text(slide, x, Inches(6.78), Inches(2.8), Inches(0.4),
             m, font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 7: תחזית פיננסית
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'תחזית פיננסית — 3 שנים', font_size=44, color=ACCENT, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.2), Inches(1.2), Inches(3), ACCENT)

# KPI boxes at top
kpis = [
    ('33.6M₪', 'ARR שנה 3', ACCENT),
    ('40,000', 'משתמשים משלמים', GREEN),
    ('66%', 'מרווח תפעולי', GOLD),
]

for i, (val, label, color) in enumerate(kpis):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.6), Inches(3.8), Inches(1.2), BG_CARD, 0.06)
    add_shape(slide, x, Inches(1.6), Inches(3.8), Inches(0.06), color)
    add_text(slide, x, Inches(1.75), Inches(3.8), Inches(0.6),
             val, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER, rtl=False)
    add_text(slide, x, Inches(2.3), Inches(3.8), Inches(0.35),
             label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Table
card = add_shape(slide, Inches(0.8), Inches(3.1), Inches(11.7), Inches(3.8), BG_CARD, 0.04)

headers = ['', 'שנה 1', 'שנה 2', 'שנה 3']
rows = [
    ('משתמשים משלמים', '2,000', '12,000', '40,000', WHITE),
    ('הכנסה חודשית', '80K₪', '650K₪', '2.8M₪', LIGHT_GRAY),
    ('הכנסה שנתית (ARR)', '960K₪', '7.8M₪', '33.6M₪', WHITE),
    ('עלויות (API + תפעול)', '65K₪/חודש', '300K₪/חודש', '950K₪/חודש', GRAY),
    ('רווח תפעולי', '15K₪/חודש', '350K₪/חודש', '1.85M₪/חודש', GREEN),
    ('מרווח תפעולי', '19%', '54%', '66%', GREEN),
]

col_widths = [Inches(3.5), Inches(2.7), Inches(2.7), Inches(2.7)]
col_starts = [Inches(1)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

for j, h in enumerate(headers):
    add_text(slide, col_starts[j], Inches(3.3), col_widths[j], Inches(0.45),
             h, font_size=16, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.RIGHT)

add_shape(slide, Inches(1.2), Inches(3.75), Inches(11), Inches(0.02), ACCENT_DIM)

for i, (label, *vals, color) in enumerate(rows):
    y = Inches(3.9 + i * 0.5)
    is_profit = 'רווח' in label or 'מרווח' in label
    add_text(slide, col_starts[0], y, col_widths[0], Inches(0.4),
             label, font_size=14, color=WHITE, bold=is_profit)
    for j, val in enumerate(vals):
        add_text(slide, col_starts[j+1], y, col_widths[j+1], Inches(0.4),
                 val, font_size=14, color=color, bold=is_profit,
                 alignment=PP_ALIGN.CENTER, rtl=False)

add_text(slide, Inches(0.8), Inches(7), Inches(11.5), Inches(0.4),
         'הנחות: המרה 3% מ-Free ל-Pro  |  Churn חודשי 5%  |  עלות API ממוצעת ~$0.03/שיחה',
         font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 8: טכנולוגיה — with tech image
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_image_with_overlay(slide, f'{ASSETS}/tech-bg.jpg',
                       Inches(0), Inches(0), Inches(13.333), Inches(7.5), overlay_alpha=80)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'סטאק טכנולוגי', font_size=44, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(1.2), Inches(2.3), GOLD)

tech_layers = [
    ('Frontend', ['Next.js 15 + React 19', 'RTL Hebrew UI', 'SSE Streaming', 'Responsive Design'], '🖥️', ACCENT),
    ('Backend', ['Next.js API Routes', 'Claude API (Sonnet 4)', 'Tool Calling (Agentic AI)', 'Prompt Caching'], '⚡', GREEN),
    ('Data', ['Yahoo Finance API', 'Real-time Quotes', 'Fundamentals & Earnings', 'Sector & ETF Analysis'], '📈', GOLD),
    ('Infra', ['Vercel (Hosting)', 'Edge Functions', 'GitHub CI/CD', 'Environment Secrets'], '☁️', ACCENT_LIGHT),
]

for i, (title, items, icon, color) in enumerate(tech_layers):
    x = Inches(0.4 + i * 3.2)
    card = add_shape(slide, x, Inches(1.8), Inches(3), Inches(5.2), BG_CARD, 0.06, alpha=92)

    # Top color bar
    add_shape(slide, x, Inches(1.8), Inches(3), Inches(0.06), color)

    add_text(slide, x, Inches(2.05), Inches(3), Inches(0.5),
             icon, font_size=36, alignment=PP_ALIGN.CENTER, rtl=False)
    add_text(slide, x, Inches(2.6), Inches(3), Inches(0.5),
             title, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER,
             rtl=False)

    # Separator
    add_shape(slide, x + Inches(0.5), Inches(3.15), Inches(2), Inches(0.02), color)

    for j, item in enumerate(items):
        item_card = add_shape(slide, x + Inches(0.2), Inches(3.4 + j * 0.7), Inches(2.6), Inches(0.55),
                              BG_CARD_ALT, 0.08)
        add_text(slide, x + Inches(0.4), Inches(3.45 + j * 0.7), Inches(2.2), Inches(0.45),
                 item, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 rtl=False)


# ═══════════════════════════════════════════════════════
# SLIDE 9: יתרון תחרותי
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'יתרון תחרותי', font_size=44, color=ACCENT, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(1.2), Inches(2.3), ACCENT)

comp_headers = ['', 'נועם AI', 'בנקים', 'Robo-Advisors', 'יועץ אנושי']
comp_rows = [
    ['עברית טבעית', '✅', '❌', '❌', '✅'],
    ['התאמה אישית אמיתית', '✅', '❌', '⚠️', '✅'],
    ['נתונים בזמן אמת', '✅', '⚠️', '✅', '⚠️'],
    ['עלות נמוכה', '✅', '❌', '✅', '❌'],
    ['זמינות 24/7', '✅', '❌', '✅', '❌'],
    ['הסבר וחינוך פיננסי', '✅', '❌', '❌', '⚠️'],
    ['בחירת מניות בודדות', '✅', '⚠️', '❌', '✅'],
]

card = add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5), BG_CARD, 0.04)

# Highlight column for Noam
highlight_col = add_shape(slide, Inches(3.75), Inches(1.6), Inches(2.2), Inches(5.5), ACCENT_DIM, 0.03, alpha=30)

comp_col_w = [Inches(2.9), Inches(2.2), Inches(2.2), Inches(2.3), Inches(2.1)]
comp_col_x = [Inches(1)]
for w in comp_col_w[:-1]:
    comp_col_x.append(comp_col_x[-1] + w)

for j, h in enumerate(comp_headers):
    color = ACCENT_LIGHT if j == 1 else GOLD
    add_text(slide, comp_col_x[j], Inches(1.8), comp_col_w[j], Inches(0.5),
             h, font_size=17, color=color, bold=True,
             alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.RIGHT)

add_shape(slide, Inches(1.2), Inches(2.3), Inches(11), Inches(0.02), ACCENT_DIM)

for i, row in enumerate(comp_rows):
    y = Inches(2.5 + i * 0.6)
    # Alternating row bg
    if i % 2 == 0:
        add_shape(slide, Inches(1), y - Inches(0.05), Inches(11.3), Inches(0.55),
                  BG_CARD_ALT, 0.02)
    for j, cell in enumerate(row):
        c = WHITE if j == 0 else (GREEN if '✅' in cell else (RED if '❌' in cell else GOLD))
        fs = 14 if j == 0 else 20
        add_text(slide, comp_col_x[j], y, comp_col_w[j], Inches(0.4),
                 cell, font_size=fs, color=c,
                 alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.RIGHT,
                 rtl=(j == 0))


# ═══════════════════════════════════════════════════════
# SLIDE 10: Roadmap
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'מפת דרכים', font_size=44, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.8), Inches(1.2), Inches(1.8), GOLD)

# Timeline bar
add_shape(slide, Inches(1.5), Inches(3.4), Inches(10.5), Inches(0.06), ACCENT_DIM)

phases = [
    ('Q3 2026', 'MVP', ['צ׳אט יועץ בסיסי', 'בניית תיק מנייתי', 'השקה ל-Beta', 'גיוס Pre-Seed'], ACCENT),
    ('Q4 2026', 'Growth', ['תמיכה באג"ח + מדדים', 'התראות שוק', 'אפליקציית מובייל', 'שיווק דיגיטלי'], GREEN),
    ('H1 2027', 'Scale', ['מעקב תיק חי', 'חיבור לברוקר', 'רגולציה (רשות ני"ע)', 'גיוס Seed ($2M)'], GOLD),
    ('H2 2027+', 'Expand', ['שוק ת"א', 'קריפטו', 'B2B לבנקים', 'שווקים נוספים'], ACCENT_LIGHT),
]

for i, (period, title, items, color) in enumerate(phases):
    x = Inches(0.4 + i * 3.2)

    # Diamond on timeline
    diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x + Inches(1.2), Inches(3.2), Inches(0.4), Inches(0.4))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = color
    diamond.line.fill.background()

    # Period label
    period_card = add_shape(slide, x + Inches(0.5), Inches(1.8), Inches(1.8), Inches(0.45), color, 0.1)
    add_text(slide, x + Inches(0.5), Inches(1.82), Inches(1.8), Inches(0.4),
             period, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             rtl=False)

    add_text(slide, x, Inches(2.45), Inches(2.8), Inches(0.45),
             title, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             rtl=False)

    # Items card
    card = add_shape(slide, x + Inches(0.05), Inches(3.9), Inches(2.7), Inches(3.1), BG_CARD, 0.06)
    add_shape(slide, x + Inches(0.05), Inches(3.9), Inches(2.7), Inches(0.05), color)

    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.25), Inches(4.15 + j * 0.6), Inches(2.3), Inches(0.45),
                 f'▸  {item}', font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 11: צוות — with team image
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_image_with_overlay(slide, f'{ASSETS}/team.jpg',
                       Inches(0), Inches(0), Inches(13.333), Inches(7.5), overlay_alpha=78)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
         'הצוות', font_size=44, color=ACCENT, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(6), Inches(1.2), Inches(1.3), ACCENT)

team = [
    ('👨‍💻', 'מייסד & CTO', 'פיתוח מוצר ו-AI', ['ארכיטקטורת מערכות AI', 'ניתוח פיננסי אלגוריתמי', 'Full-stack development'], ACCENT),
    ('📈', 'שותף — פיננסים', '[למלא]', ['רקע בשוק ההון', 'רגולציה ופיננסים', 'ניהול תיקי השקעות'], GOLD),
    ('🎨', 'שותף — מוצר', '[למלא]', ['UX/UI וחוויית משתמש', 'שיווק דיגיטלי', 'Growth hacking'], GREEN),
]

for i, (icon, title, name, skills, color) in enumerate(team):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(5), BG_CARD, 0.06, alpha=92)
    add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(0.06), color)

    # Avatar circle
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), Inches(2.1), Inches(1.2), Inches(1.2))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = BG_CARD_ALT
    avatar.line.color.rgb = color
    avatar.line.width = Pt(2)
    add_text(slide, x + Inches(1.3), Inches(2.2), Inches(1.2), Inches(1),
             icon, font_size=44, alignment=PP_ALIGN.CENTER)

    add_text(slide, x, Inches(3.4), Inches(3.8), Inches(0.45),
             title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.85), Inches(3.8), Inches(0.35),
             name, font_size=15, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Separator
    add_shape(slide, x + Inches(0.8), Inches(4.3), Inches(2.2), Inches(0.02), color)

    for j, skill in enumerate(skills):
        add_text(slide, x + Inches(0.3), Inches(4.5 + j * 0.5), Inches(3.2), Inches(0.4),
                 f'▸  {skill}', font_size=13, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 12: סיכום ו-CTA — with finance bg
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_image_with_overlay(slide, f'{ASSETS}/finance-bg.jpg',
                       Inches(0), Inches(0), Inches(13.333), Inches(7.5), overlay_alpha=78)

# Decorative circles
add_glow_circle(slide, Inches(2), Inches(6), Inches(0.5), ACCENT_DIM)
add_glow_circle(slide, Inches(11), Inches(1.5), Inches(0.35), GOLD_DIM)

add_text(slide, Inches(1), Inches(0.8), Inches(11.333), Inches(1),
         'נועם — ייעוץ השקעות לכולם', font_size=52, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(1.9), Inches(3.333), ACCENT)

highlights = [
    ('🎯', 'שוק של 2B₪ בישראל — אפס מתחרים ב-AI בעברית'),
    ('💰', 'מודל SaaS רווחי — מרווח 66% בשנה 3'),
    ('🚀', 'MVP עובד וחי באוויר — מוכן לצמיחה'),
    ('🔒', 'First Mover — יתרון ראשון לשוק'),
]

for i, (icon, text) in enumerate(highlights):
    y = Inches(2.3 + i * 0.75)
    card = add_shape(slide, Inches(2.5), y, Inches(8.333), Inches(0.6), BG_CARD, 0.08, alpha=88)
    add_text(slide, Inches(2.7), y + Inches(0.05), Inches(0.6), Inches(0.5),
             icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.3), y + Inches(0.08), Inches(7), Inches(0.45),
             text, font_size=19, color=LIGHT_GRAY)

# iPhone mockups flanking CTA
iphone_welcome = f'{ASSETS}/iphone-app-screenshot.png'
iphone_chat_img = f'{ASSETS}/iphone-screen-chat.png'
if os.path.exists(iphone_welcome):
    slide.shapes.add_picture(iphone_welcome, Inches(0.3), Inches(2), Inches(2.2))
if os.path.exists(iphone_chat_img):
    slide.shapes.add_picture(iphone_chat_img, Inches(10.8), Inches(2), Inches(2.2))

add_shape(slide, Inches(4), Inches(5.5), Inches(5.333), Inches(1.1), ACCENT, 0.12)
add_text(slide, Inches(4), Inches(5.55), Inches(5.333), Inches(0.55),
         'בואו נבנה את העתיד של ייעוץ ההשקעות', font_size=22, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4), Inches(6.05), Inches(5.333), Inches(0.4),
         'investment-chat-psi.vercel.app', font_size=16, color=ACCENT_LIGHT,
         alignment=PP_ALIGN.CENTER, rtl=False)

add_text(slide, Inches(1), Inches(6.9), Inches(11.333), Inches(0.4),
         'shaharezerzer@gmail.com', font_size=16, color=GRAY,
         alignment=PP_ALIGN.CENTER, rtl=False)


# ── Save ──
output_path = os.path.expanduser('~/Desktop/נועם_יועץ_השקעות_AI_תכנון_עסקי.pptx')
prs.save(output_path)
print(f'Saved: {output_path}')
